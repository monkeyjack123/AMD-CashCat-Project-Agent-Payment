from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path(__file__).resolve().parent / "CashCat-AMD-Track1-Deck.pptx"

BG = RGBColor(7, 17, 26)
BG_2 = RGBColor(11, 24, 36)
PANEL = RGBColor(14, 26, 42)
PANEL_SOFT = RGBColor(20, 37, 55)
TEXT = RGBColor(238, 248, 255)
MUTED = RGBColor(145, 168, 188)
CYAN = RGBColor(87, 236, 255)
GREEN = RGBColor(115, 255, 186)
BLUE = RGBColor(122, 159, 255)
AMBER = RGBColor(255, 209, 102)
RED = RGBColor(255, 110, 141)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    top_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15)
    )
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = BG_2
    top_band.line.fill.background()

    glow_left = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.8), Inches(-0.4), Inches(3.3), Inches(2.2)
    )
    glow_left.fill.solid()
    glow_left.fill.fore_color.rgb = CYAN
    glow_left.fill.transparency = 0.82
    glow_left.line.fill.background()

    glow_right = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.1), Inches(0.1), Inches(3.0), Inches(2.0)
    )
    glow_right.fill.solid()
    glow_right.fill.fore_color.rgb = BLUE
    glow_right.fill.transparency = 0.86
    glow_right.line.fill.background()

    footer = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "CashCat  |  AMD Developer Hackathon  |  Track 1: AI Agents & Agentic Workflows"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def add_brand(slide, page_num):
    pill = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.42), Inches(2.35), Inches(0.35)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = PANEL_SOFT
    pill.line.color.rgb = CYAN
    pill.line.transparency = 0.55

    tf = pill.text_frame
    p = tf.paragraphs[0]
    p.text = "CASHCAT"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CYAN

    num = slide.shapes.add_textbox(Inches(12.1), Inches(0.42), Inches(0.65), Inches(0.3))
    p = num.text_frame.paragraphs[0]
    p.text = f"{page_num:02d}"
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = MUTED


def add_title(slide, eyebrow, title, subtitle=None):
    brow = slide.shapes.add_textbox(Inches(0.72), Inches(1.05), Inches(3.0), Inches(0.3))
    p = brow.text_frame.paragraphs[0]
    p.text = eyebrow.upper()
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CYAN

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(8.4), Inches(1.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.72), Inches(2.28), Inches(8.4), Inches(0.75))
        tf = sub.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(15)
        p.font.color.rgb = MUTED


def add_bullet_panel(slide, x, y, w, h, heading, bullets, accent=CYAN):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = accent
    panel.line.transparency = 0.68

    header = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.16), w - Inches(0.4), Inches(0.3))
    p = header.text_frame.paragraphs[0]
    p.text = heading
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = accent

    body = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.52), w - Inches(0.35), h - Inches(0.65))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT
        p.level = 0
        p.bullet = True
        p.space_after = Pt(9)


def add_metric_card(slide, x, y, w, h, label, value, accent=GREEN):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL_SOFT
    card.line.color.rgb = accent
    card.line.transparency = 0.74

    label_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.16), w - Inches(0.3), Inches(0.3))
    p = label_box.text_frame.paragraphs[0]
    p.text = label.upper()
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = MUTED

    val_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.45), w - Inches(0.3), h - Inches(0.55))
    tf = val_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT


def add_flow_boxes(slide, labels):
    x = 0.8
    y = 3.15
    w = 2.35
    h = 1.2
    gap = 0.42
    colors = [CYAN, BLUE, GREEN, AMBER]
    for i, label in enumerate(labels):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + i * (w + gap)), Inches(y), Inches(w), Inches(h)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = PANEL
        box.line.color.rgb = colors[i % len(colors)]
        box.line.transparency = 0.45

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = TEXT

        if i < len(labels) - 1:
            arrow = slide.shapes.add_textbox(
                Inches(x + i * (w + gap) + w + 0.08), Inches(y + 0.38), Inches(0.26), Inches(0.3)
            )
            p = arrow.text_frame.paragraphs[0]
            p.text = "→"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = MUTED


def add_quote(slide, text):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(5.65), Inches(11.65), Inches(0.9)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = PANEL_SOFT
    box.line.color.rgb = CYAN
    box.line.transparency = 0.76

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT
    p.alignment = PP_ALIGN.CENTER


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_brand(slide, 1)
    add_title(
        slide,
        "Problem",
        "AI agents can decide what to do, but they still cannot spend safely.",
        "LLMs can plan, evaluate vendors, and trigger actions. The missing layer is governed money movement.",
    )
    add_bullet_panel(
        slide,
        Inches(0.72),
        Inches(3.05),
        Inches(5.8),
        Inches(2.2),
        "What breaks today",
        [
            "Giving an agent raw card or wallet access is unsafe.",
            "Finance teams need limits, approvals, and evidence.",
            "Most AI workflows stop right before the real economic action.",
        ],
        RED,
    )
    add_bullet_panel(
        slide,
        Inches(6.8),
        Inches(3.05),
        Inches(5.8),
        Inches(2.2),
        "What should happen instead",
        [
            "Agents should get bounded spend authority, not unrestricted money.",
            "Every delegated purchase should be explainable and auditable.",
            "AI should be able to complete workflows, not just recommend them.",
        ],
        CYAN,
    )
    add_quote(slide, "The gap is not intelligence. The gap is safe execution.")


def slide_why_now(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_brand(slide, 2)
    add_title(
        slide,
        "Why Now",
        "Agent workflows are moving from assistants to operators.",
        "The timing is driven by open models, lower inference costs, and growing willingness to let software act on behalf of teams.",
    )
    add_metric_card(slide, Inches(0.82), Inches(3.15), Inches(2.9), Inches(1.4), "Shift 1", "Open models can now plan and structure spend decisions", CYAN)
    add_metric_card(slide, Inches(3.95), Inches(3.15), Inches(2.9), Inches(1.4), "Shift 2", "AI-native teams are buying APIs, data, SaaS, and cloud on the fly", BLUE)
    add_metric_card(slide, Inches(7.08), Inches(3.15), Inches(2.9), Inches(1.4), "Shift 3", "Enterprises still need approvals, controls, and audit trails", GREEN)
    add_metric_card(slide, Inches(10.2), Inches(3.15), Inches(2.1), Inches(1.4), "Result", "New control layer", AMBER)
    add_bullet_panel(
        slide,
        Inches(0.82),
        Inches(5.0),
        Inches(11.45),
        Inches(1.35),
        "Category opening",
        [
            "As soon as agents can buy digital resources, the market shifts from copilots to AI-native operations. CashCat sits exactly at that transition."
        ],
        CYAN,
    )


def slide_product(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_brand(slide, 3)
    add_title(
        slide,
        "Product",
        "CashCat is the payment and spend control layer for AI agents.",
        "It lets agents buy APIs, data, software, and cloud resources under human-defined rules.",
    )
    add_bullet_panel(
        slide,
        Inches(0.82),
        Inches(3.1),
        Inches(4.0),
        Inches(2.4),
        "AI does",
        [
            "Interprets the task",
            "Selects paid resources",
            "Generates purchase proposals",
            "Triggers workflow execution",
        ],
        CYAN,
    )
    add_bullet_panel(
        slide,
        Inches(4.98),
        Inches(3.1),
        Inches(3.2),
        Inches(2.4),
        "CashCat does",
        [
            "Budget limits",
            "Category restrictions",
            "Approval thresholds",
            "Receipts and spend proofs",
        ],
        GREEN,
    )
    add_bullet_panel(
        slide,
        Inches(8.36),
        Inches(3.1),
        Inches(3.95),
        Inches(2.4),
        "Buyer value",
        [
            "No raw credential exposure",
            "Safer delegated spend",
            "Clear auditability",
            "AI can complete real work",
        ],
        BLUE,
    )
    add_quote(slide, "The AI decides. CashCat governs.")


def slide_how_it_works(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_brand(slide, 4)
    add_title(
        slide,
        "How It Works",
        "From natural language task to governed economic action.",
        "The system keeps the user experience simple while activating orchestration when the task becomes more complex.",
    )
    add_flow_boxes(
        slide,
        [
            "User gives the agent a business task",
            "Qwen on AMD proposes the paid actions",
            "CashCat applies budget, approval, and policy",
            "Workflow artifacts, receipt, and proof are returned",
        ],
    )
    add_bullet_panel(
        slide,
        Inches(0.82),
        Inches(4.9),
        Inches(11.45),
        Inches(1.3),
        "Behind the scenes",
        [
            "Simple tasks stay single-agent. Broader workflows can automatically escalate into orchestrated specialist lanes without changing the user entry point."
        ],
        AMBER,
    )


def slide_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_brand(slide, 5)
    add_title(
        slide,
        "Live Demo Flow",
        "One prompt. One entry point. Real AI evidence and governed outputs.",
        "The demo shows model-backed planning, spend governance, and visible workflow artifacts rather than a static front-end mock.",
    )
    add_bullet_panel(
        slide,
        Inches(0.82),
        Inches(3.05),
        Inches(3.5),
        Inches(2.45),
        "User sees",
        [
            "Single chat-style task box",
            "AI thinking and staged execution",
            "A proposal, a control decision, and a final artifact",
        ],
        CYAN,
    )
    add_bullet_panel(
        slide,
        Inches(4.53),
        Inches(3.05),
        Inches(3.8),
        Inches(2.45),
        "System proves",
        [
            "Model: Qwen2.5-14B-Instruct",
            "Provider: AMD vLLM endpoint",
            "Payment intent, receipt, and spend proof",
        ],
        GREEN,
    )
    add_bullet_panel(
        slide,
        Inches(8.55),
        Inches(3.05),
        Inches(3.7),
        Inches(2.45),
        "Workflow results",
        [
            "Research brief",
            "Benchmark result",
            "Vendor stack or launch plan",
        ],
        BLUE,
    )
    add_quote(slide, "The product does not stop at recommendation. It pushes all the way to governed action and a visible deliverable.")


def slide_market(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_brand(slide, 6)
    add_title(
        slide,
        "Market Opportunity",
        "We start with digital spend, but the category expands with every AI-triggered purchase.",
        "The wedge is narrow enough to ship today and broad enough to become a platform.",
    )
    add_bullet_panel(
        slide,
        Inches(0.82),
        Inches(3.08),
        Inches(3.6),
        Inches(2.45),
        "Initial wedge",
        [
            "APIs",
            "Datasets",
            "SaaS tools",
            "Cloud resources",
        ],
        CYAN,
    )
    add_bullet_panel(
        slide,
        Inches(4.58),
        Inches(3.08),
        Inches(3.7),
        Inches(2.45),
        "Expansion",
        [
            "AI procurement",
            "Delegated operational spend",
            "Agent-to-agent commerce",
            "Enterprise finance controls",
        ],
        GREEN,
    )
    add_bullet_panel(
        slide,
        Inches(8.46),
        Inches(3.08),
        Inches(3.8),
        Inches(2.45),
        "Why it can be large",
        [
            "Every useful agent eventually needs tools, data, services, or compute.",
            "Wherever AI touches budget, CashCat can become the control plane.",
        ],
        BLUE,
    )
    add_quote(slide, "We are not just selling payments. We are defining the financial operating layer for AI agents.")


def slide_business(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_brand(slide, 7)
    add_title(
        slide,
        "Business Model",
        "CashCat monetizes like an AI-native blend of spend control software and payments infrastructure.",
        "The model starts simple and gets stronger as agent-driven volume grows.",
    )
    add_metric_card(slide, Inches(0.82), Inches(3.15), Inches(3.6), Inches(1.55), "Revenue stream 1", "Platform subscription for teams using AI spend controls", CYAN)
    add_metric_card(slide, Inches(4.58), Inches(3.15), Inches(3.6), Inches(1.55), "Revenue stream 2", "Usage or API pricing for agent builders and workflow platforms", GREEN)
    add_metric_card(slide, Inches(8.34), Inches(3.15), Inches(3.92), Inches(1.55), "Revenue stream 3", "Transaction or spend-based monetization as governed volume scales", BLUE)
    add_bullet_panel(
        slide,
        Inches(0.82),
        Inches(5.12),
        Inches(11.42),
        Inches(1.15),
        "Go-to-market logic",
        [
            "Start with AI-native teams, agent builders, and innovation groups. Expand into enterprise controls, compliance workflows, and higher-volume delegated spend."
        ],
        AMBER,
    )


def slide_vision_amd(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_brand(slide, 8)
    add_title(
        slide,
        "Vision + Why AMD",
        "CashCat turns AI from a recommender into an operator with governed economic power.",
        "AMD cloud lets us pair open models with private, controllable, high-performance inference for agent decisioning.",
    )
    add_bullet_panel(
        slide,
        Inches(0.82),
        Inches(3.0),
        Inches(5.35),
        Inches(2.65),
        "Why AMD matters",
        [
            "Open-model inference on dedicated GPU infrastructure",
            "Control over latency, privacy, and deployment",
            "A credible path from hackathon demo to production architecture",
        ],
        CYAN,
    )
    add_bullet_panel(
        slide,
        Inches(6.42),
        Inches(3.0),
        Inches(5.82),
        Inches(2.65),
        "Long-term vision",
        [
            "Every serious AI agent will need governed money movement.",
            "CashCat can become the standard way agents receive spend authority, execute safely, and return proof.",
            "That is a platform, not a feature.",
        ],
        GREEN,
    )
    add_quote(slide, "If AI is going to act in the economy, it needs a financial control layer. That layer is CashCat.")


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_problem(prs)
    slide_why_now(prs)
    slide_product(prs)
    slide_how_it_works(prs)
    slide_demo(prs)
    slide_market(prs)
    slide_business(prs)
    slide_vision_amd(prs)

    prs.save(str(OUT_PATH))
    print(OUT_PATH)


if __name__ == "__main__":
    build()
